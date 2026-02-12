"""
PPTX Exporter — generates PowerPoint files from slide data.

Supports two modes:
  1. From scratch: Creates a .pptx with programmatic styling based on preset colors/fonts
  2. Template-based: Uses a user-provided .pptx as a template, injecting content into existing layouts
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger("slide-builder.exporters.pptx")


def export_pptx(
    title: str,
    slides: list[dict[str, Any]],
    style_name: str,
    output_path: str,
    custom_preset: dict[str, Any] | None = None,
    template_path: str | None = None,
) -> str:
    """
    Generate a PPTX file from slide data.

    Args:
        title: Presentation title.
        slides: List of slide dicts.
        style_name: Preset name for styling.
        output_path: Where to write the .pptx file.
        custom_preset: Optional custom preset dict (overrides style_name).
        template_path: Optional .pptx template file to use as base.

    Returns:
        Absolute path to the generated .pptx file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx required: pip install python-pptx")

    # Load preset colors/fonts
    preset = _load_style(style_name, custom_preset)
    colors = preset.get("colors", {})
    fonts = preset.get("fonts", {})

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Set 16:9 aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Parse colors
    bg_color = _hex_to_rgb(colors.get("bg_primary", "#1a1a2e"))
    text_color = _hex_to_rgb(colors.get("text_primary", "#ffffff"))
    text_secondary = _hex_to_rgb(colors.get("text_secondary", "#b0b0b0"))
    accent_color = _hex_to_rgb(colors.get("accent", "#e94560"))

    display_font = fonts.get("display", {}).get("family", "Arial")
    body_font = fonts.get("body", {}).get("family", "Arial")

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")

        if template_path and Path(template_path).exists():
            # Template mode: find best matching layout
            layout = _find_best_layout(prs, slide_type)
        else:
            # From scratch: use blank layout
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Build slide content based on type
        if slide_type == "title":
            _build_title_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )
        elif slide_type == "quote":
            _build_quote_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )
        elif slide_type == "feature_grid":
            _build_feature_grid_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )
        elif slide_type == "code":
            _build_code_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )
        elif slide_type == "closing":
            _build_closing_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )
        else:
            # Default: content slide with bullets
            _build_content_slide(
                slide, slide_data, display_font, body_font,
                text_color, text_secondary, accent_color,
            )

        # Add speaker notes
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

    # Save
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out.resolve())


def _load_style(style_name: str, custom_preset: dict[str, Any] | None) -> dict[str, Any]:
    """Load style preset."""
    if custom_preset:
        return custom_preset
    try:
        from ..styles import load_preset
        return load_preset(style_name)
    except Exception:
        return {
            "colors": {
                "bg_primary": "#1a1a2e",
                "text_primary": "#ffffff",
                "text_secondary": "#b0b0b0",
                "accent": "#e94560",
            },
            "fonts": {
                "display": {"family": "Arial"},
                "body": {"family": "Arial"},
            },
        }


def _hex_to_rgb(hex_color: str) -> Any:
    """Convert hex color string to RGBColor."""
    from pptx.dml.color import RGBColor

    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "1a1a2e"
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _find_best_layout(prs: Any, slide_type: str) -> Any:
    """Find the best matching layout in a template for a slide type."""
    type_to_layout_keywords = {
        "title": ["title slide", "title", "cover"],
        "content": ["title and content", "content", "two content"],
        "feature_grid": ["title and content", "two content", "content"],
        "quote": ["blank", "title only"],
        "code": ["blank", "title only"],
        "closing": ["title slide", "title", "blank"],
        "image": ["picture", "title and content", "blank"],
    }

    keywords = type_to_layout_keywords.get(slide_type, ["blank"])

    for keyword in keywords:
        for layout in prs.slide_layouts:
            if keyword.lower() in layout.name.lower():
                return layout

    # Fallback: use blank or first layout
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[0]


def _add_textbox(
    slide: Any,
    left: float, top: float, width: float, height: float,
    text: str, font_name: str, font_size: int,
    font_color: Any, bold: bool = False, alignment: int = 0,
) -> Any:
    """Add a text box to a slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold

    align_map = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT}
    p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

    return txBox


def _build_title_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a title slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # Add accent bar at top
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(0.1)  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    # Title
    _add_textbox(
        slide, 1.0, 2.5, 11.333, 2.0,
        data.get("title", ""), display_font, 44,
        text_color, bold=True, alignment=1,
    )

    # Subtitle
    if data.get("subtitle"):
        _add_textbox(
            slide, 1.5, 4.5, 10.333, 1.5,
            data["subtitle"], body_font, 24,
            text_secondary, alignment=1,
        )


def _build_content_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a content slide with bullets."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # Title
    _add_textbox(
        slide, 0.8, 0.5, 11.733, 1.0,
        data.get("title", ""), display_font, 32,
        text_color, bold=True,
    )

    # Accent underline
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.4), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    # Bullets
    bullets = data.get("bullets", [])
    if bullets:
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"  •  {bullet}"
            p.font.name = body_font
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(12)


def _build_quote_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a quote slide."""
    from pptx.util import Inches, Pt

    # Large quote mark
    _add_textbox(
        slide, 1.5, 1.0, 1.5, 1.5,
        "\u201C", display_font, 96,
        accent_color, alignment=1,
    )

    # Quote text
    _add_textbox(
        slide, 2.0, 2.5, 9.333, 3.0,
        data.get("quote", ""), display_font, 28,
        text_color, alignment=1,
    )

    # Attribution
    if data.get("attribution"):
        _add_textbox(
            slide, 2.0, 5.5, 9.333, 1.0,
            f"\u2014 {data['attribution']}", body_font, 18,
            text_secondary, alignment=1,
        )


def _build_feature_grid_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a feature grid slide with cards."""
    from pptx.util import Inches, Pt

    # Title
    _add_textbox(
        slide, 0.8, 0.5, 11.733, 1.0,
        data.get("title", ""), display_font, 32,
        text_color, bold=True,
    )

    cards = data.get("cards", [])[:6]
    if not cards:
        return

    # Grid layout: 2 rows x 3 cols or 1 row x 3 cols
    cols = min(3, len(cards))
    rows = (len(cards) + cols - 1) // cols
    card_width = 3.5
    card_height = 2.2
    start_x = (13.333 - (cols * card_width + (cols - 1) * 0.3)) / 2
    start_y = 1.8

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_width + 0.3)
        y = start_y + row * (card_height + 0.3)

        # Card background
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(card_width), Inches(card_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb("#ffffff0d"[:7] if len("#ffffff0d") > 7 else "#2a2a3e")
        shape.line.fill.background()

        # Card icon
        icon = card.get("icon", "")
        if icon:
            _add_textbox(
                slide, x + 0.2, y + 0.2, 0.5, 0.5,
                icon, body_font, 20, accent_color,
            )

        # Card title
        _add_textbox(
            slide, x + 0.2, y + 0.7, card_width - 0.4, 0.5,
            card.get("title", ""), body_font, 16,
            text_color, bold=True,
        )

        # Card description
        _add_textbox(
            slide, x + 0.2, y + 1.2, card_width - 0.4, 0.9,
            card.get("description", ""), body_font, 12,
            text_secondary,
        )


def _build_code_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a code slide."""
    from pptx.util import Inches, Pt

    # Title
    _add_textbox(
        slide, 0.8, 0.5, 11.733, 1.0,
        data.get("title", ""), display_font, 32,
        text_color, bold=True,
    )

    # Code block background
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb("#0d0d0d")
    shape.line.fill.background()

    # Code text
    code = data.get("code", "")
    _add_textbox(
        slide, 1.0, 2.0, 11.333, 4.5,
        code, "Courier New", 14,
        _hex_to_rgb("#00ff00"),
    )


def _build_closing_slide(
    slide: Any, data: dict, display_font: str, body_font: str,
    text_color: Any, text_secondary: Any, accent_color: Any,
) -> None:
    """Build a closing slide."""
    from pptx.util import Inches, Pt

    # Accent bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    # Title
    _add_textbox(
        slide, 1.0, 2.5, 11.333, 2.0,
        data.get("title", "Thank You"), display_font, 44,
        text_color, bold=True, alignment=1,
    )

    # Subtitle
    if data.get("subtitle"):
        _add_textbox(
            slide, 1.5, 4.5, 10.333, 1.5,
            data["subtitle"], body_font, 24,
            text_secondary, alignment=1,
        )
