"""
Slide Builder MCP — PowerPoint Converter.

Extracts content from .pptx files using python-pptx
and converts to structured slide data for HTML generation.
"""

from __future__ import annotations

import base64
import json
import os
from pathlib import Path
from typing import Any


def extract_pptx(file_path: str, output_dir: str) -> dict[str, Any]:
    """
    Extract all content from a PowerPoint file.

    Returns a JSON-serializable structure with slides, text, and images.
    Images are saved to {output_dir}/assets/ and referenced by path.

    Args:
        file_path: Path to the .pptx file.
        output_dir: Directory for output (assets saved here).

    Returns:
        Dict with 'slides' list, 'total_slides' count, and 'assets_dir' path.
    """
    # Import here so we get a clear error if python-pptx is not installed
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPT conversion. "
            "Install it with: pip install python-pptx"
        )

    prs = Presentation(file_path)
    slides_data: list[dict[str, Any]] = []

    # Create assets directory
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data: dict[str, Any] = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
            "layout": _detect_layout(slide),
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data["title"] = shape.text.strip()
                else:
                    # Extract all text runs with formatting
                    paragraphs = []
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            para_data: dict[str, Any] = {"text": text}
                            # Detect if bullet
                            if paragraph.level > 0:
                                para_data["level"] = paragraph.level
                                para_data["type"] = "bullet"
                            paragraphs.append(para_data)
                    if paragraphs:
                        slide_data["content"].append({
                            "type": "text",
                            "paragraphs": paragraphs,
                        })

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_filename = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}"

                # Determine extension from content type
                ext_map = {
                    "image/png": ".png",
                    "image/jpeg": ".jpg",
                    "image/gif": ".gif",
                    "image/svg+xml": ".svg",
                    "image/bmp": ".bmp",
                    "image/tiff": ".tiff",
                }
                ext = ext_map.get(image.content_type, ".png")
                image_filename += ext

                # Save image
                image_path = os.path.join(assets_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image.blob)

                slide_data["images"].append({
                    "filename": image_filename,
                    "path": f"assets/{image_filename}",
                    "content_type": image.content_type,
                    "width": shape.width,
                    "height": shape.height,
                })

            # Extract tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                slide_data["content"].append({
                    "type": "table",
                    "rows": table_data,
                })

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append(slide_data)

    return {
        "slides": slides_data,
        "total_slides": len(slides_data),
        "assets_dir": assets_dir,
        "source_file": os.path.basename(file_path),
    }


def _detect_layout(slide: Any) -> str:
    """Attempt to detect the slide layout type."""
    layout_name = ""
    try:
        layout_name = slide.slide_layout.name.lower()
    except Exception:
        pass

    if "title" in layout_name and "only" in layout_name:
        return "title"
    if "title" in layout_name:
        return "title"
    if "blank" in layout_name:
        return "blank"
    if "two" in layout_name or "comparison" in layout_name:
        return "two_column"
    if "content" in layout_name:
        return "content"
    return "content"


def pptx_to_slides(file_path: str, output_dir: str) -> list[dict[str, Any]]:
    """
    Convert extracted PPTX data to the slide format expected by the generator.

    Returns a list of slide dicts compatible with generator.generate_presentation().
    """
    extracted = extract_pptx(file_path, output_dir)
    slides: list[dict[str, Any]] = []

    for i, ppt_slide in enumerate(extracted["slides"]):
        slide: dict[str, Any] = {"title": ppt_slide["title"]}

        # First slide → title type
        if i == 0:
            slide["type"] = "title"
            # Use first content paragraph as subtitle
            for content in ppt_slide["content"]:
                if content["type"] == "text" and content["paragraphs"]:
                    slide["subtitle"] = content["paragraphs"][0]["text"]
                    break
        # Last slide with short content → closing
        elif i == len(extracted["slides"]) - 1 and len(ppt_slide["content"]) <= 1:
            slide["type"] = "closing"
            for content in ppt_slide["content"]:
                if content["type"] == "text" and content["paragraphs"]:
                    slide["subtitle"] = content["paragraphs"][0]["text"]
                    break
        # Slide with images → image type
        elif ppt_slide["images"]:
            slide["type"] = "image"
            slide["image_src"] = ppt_slide["images"][0]["path"]
            # Still include text content as bullets
            bullets = []
            for content in ppt_slide["content"]:
                if content["type"] == "text":
                    for para in content["paragraphs"]:
                        bullets.append(para["text"])
            if bullets:
                slide["bullets"] = bullets[:6]
        # Default → content with bullets
        else:
            slide["type"] = "content"
            bullets = []
            for content in ppt_slide["content"]:
                if content["type"] == "text":
                    for para in content["paragraphs"]:
                        bullets.append(para["text"])
            if bullets:
                slide["bullets"] = bullets[:6]

        slides.append(slide)

    return slides


def summarize_extraction(file_path: str, output_dir: str) -> str:
    """
    Extract a PPTX and return a human-readable summary for user confirmation.
    """
    extracted = extract_pptx(file_path, output_dir)
    lines = [
        f"Extracted **{extracted['total_slides']}** slides from `{extracted['source_file']}`:\n"
    ]

    for slide in extracted["slides"]:
        img_count = len(slide["images"])
        content_count = sum(
            len(c.get("paragraphs", [])) for c in slide["content"] if c["type"] == "text"
        )
        img_info = f" | {img_count} image(s)" if img_count else ""
        title = slide["title"] or "(No title)"
        lines.append(f"**Slide {slide['number']}: {title}**")
        lines.append(f"  - {content_count} text block(s){img_info}")
        if slide["notes"]:
            lines.append(f"  - Speaker notes: {slide['notes'][:80]}...")
        lines.append("")

    lines.append(f"\nAll images saved to `{extracted['assets_dir']}`")
    return "\n".join(lines)
